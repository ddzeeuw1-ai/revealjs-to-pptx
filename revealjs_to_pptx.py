#!/usr/bin/env python3
"""
revealjs_to_pptx.py
────────────────────────────────────────────────────────────────────────────────
Converts a reveal.js HTML slide deck (publictimes_CMYK / base_2022 design
system) into an editable PowerPoint (.pptx) file.

Usage
─────
    python revealjs_to_pptx.py <path/to/slides.html> [--out output.pptx]

Expected structure:
    slides.html
    img/   ← background and content images

Dependencies
────────────
    pip install python-pptx beautifulsoup4

Design system (from CSS)
────────────────────────
  Background  #000000   Cyan    #00CCFF   Magenta  #FF0099
  White       #FFFFFF   Yellow  #FFEE00   CyanL    #CCCCFF
  Fonts: Public Sans → Calibri   Happy Times → Georgia

NOTE on this deck's HTML quirk
───────────────────────────────
All <ul>/<ol> lists use NESTED <li> elements (each <li> wraps the next)
instead of sibling elements.  Use flatten_li() to unwrap them.
"""

import argparse, re, sys
from pathlib import Path
from bs4 import BeautifulSoup, NavigableString, Tag
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Canvas ────────────────────────────────────────────────────────────────────
W = Inches(10)
H = Inches(5.625)

# ── Colours ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x00, 0x00, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
CYAN    = RGBColor(0x00, 0xCC, 0xFF)
MAGENTA = RGBColor(0xFF, 0x00, 0x99)
YELLOW  = RGBColor(0xFF, 0xEE, 0x00)
MUTED   = RGBColor(0x99, 0x99, 0x99)
CYAN_L  = RGBColor(0xCC, 0xCC, 0xFF)

FONT_SANS  = "Calibri"
FONT_SERIF = "Georgia"   # stands in for Happy Times


# ═══════════════════════════════════════════════════════════════════════════════
# HTML structural helpers
# ═══════════════════════════════════════════════════════════════════════════════

def flatten_li(list_tag):
    """
    This deck nests <li> inside each other instead of making them siblings.
    Walk the chain and collect each <li>'s OWN direct text/inline content,
    stopping before any nested <li> child.
    Returns a list of Tag objects (synthetic clones with only own content).
    """
    items = []
    li = list_tag.find("li") if list_tag else None
    while li:
        # Build a shallow clone with only non-li children
        clone = BeautifulSoup("<span></span>", "html.parser").find("span")
        for child in li.children:
            if isinstance(child, NavigableString):
                clone.append(type(child)(str(child)))
            elif isinstance(child, Tag) and child.name != "li":
                clone.append(child.__copy__())
        items.append(clone)
        li = li.find("li", recursive=False)
    return items


def get_col0(section):
    """Return the first (content) div inside a .twothird or .twocol container."""
    tw = section.find(class_=re.compile(r"\b(twothird|twocol)\b"))
    if not tw:
        return None
    divs = [c for c in tw.children if isinstance(c, Tag) and c.name == "div"]
    return divs[0] if divs else tw


def get_cols(section):
    """Return list of all column divs inside a .twocol container."""
    tw = section.find(class_=re.compile(r"\btwocol\b"))
    if not tw:
        return []
    return [c for c in tw.children if isinstance(c, Tag) and c.name == "div"]


def resolve_img(src, img_dir):
    return img_dir / Path(src).name if src else None


# ═══════════════════════════════════════════════════════════════════════════════
# Shape / text helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _in(n):
    return Inches(n)


def add_rect(slide, x, y, w, h, fill_rgb, transparency=0):
    """Filled rectangle; transparency 0=opaque 100=invisible."""
    from lxml import etree
    NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    shape = slide.shapes.add_shape(1, _in(x), _in(y), _in(w), _in(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if transparency:
        spPr = shape._element.spPr
        sf = spPr.find(f"{{{NS}}}solidFill")
        if sf is not None:
            srgb = sf.find(f"{{{NS}}}srgbClr")
            if srgb is not None:
                a = etree.SubElement(srgb, f"{{{NS}}}alpha")
                a.set("val", str(int((100 - transparency) * 1000)))
    return shape


def set_black_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def set_image_bg(slide, img_path, transparency=25):
    if img_path and Path(img_path).exists():
        sp = slide.shapes.add_picture(str(img_path), _in(0), _in(0), W, H)
        slide.shapes._spTree.remove(sp._element)
        slide.shapes._spTree.insert(2, sp._element)
    add_rect(slide, 0, 0, 10, 5.625, BG, transparency=transparency)


def add_image(slide, img_path, x, y, w, h):
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), _in(x), _in(y), _in(w), _in(h))
    else:
        ph = add_rect(slide, x, y, w, h, MUTED)
        ph.text_frame.text = f"[{Path(img_path).name if img_path else 'img'}]"


def txbox(slide, x, y, w, h):
    """Return a word-wrapping textbox."""
    tb = slide.shapes.add_textbox(_in(x), _in(y), _in(w), _in(h))
    tb.word_wrap = True
    tb.text_frame.word_wrap = True
    return tb


def simple_text(slide, x, y, w, h, text, font=FONT_SANS, size=14,
                bold=False, italic=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = txbox(slide, x, y, w, h)
    p = tb.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font;  r.font.size = Pt(size)
    r.font.bold = bold;  r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_running_head(slide, label):
    if not label:
        return
    text = label.upper()
    pill_w = max(len(text) * 0.092 + 0.3, 0.7)
    add_rect(slide, 0, 0, pill_w, 0.27, MAGENTA)
    tb = txbox(slide, 0.06, 0.0, pill_w, 0.27)
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text;  r.font.name = FONT_SANS;  r.font.size = Pt(8)
    r.font.bold = True;  r.font.color.rgb = WHITE


def add_counter(slide, n, total):
    simple_text(slide, 8.5, 5.33, 1.3, 0.22, f"{n} / {total}",
                size=9, color=YELLOW, align=PP_ALIGN.RIGHT)


# ── Inline HTML → pptx paragraph ─────────────────────────────────────────────

def parse_inline(node, para, size=14, color=WHITE, font=FONT_SANS,
                 bold=False, italic=False):
    """Walk inline children adding runs; <strong>→cyan, <em>→italic, <a>→yellow."""
    for child in node.children:
        if isinstance(child, NavigableString):
            t = str(child)
            if not t:
                continue
            r = para.add_run()
            r.text = t
            r.font.name = font;   r.font.size = Pt(size)
            r.font.bold = bold;   r.font.italic = italic
            r.font.color.rgb = color
        elif isinstance(child, Tag):
            n = (child.name or "").lower()
            if n in ("strong", "b"):
                parse_inline(child, para, size, CYAN,  font, bold=True,  italic=italic)
            elif n in ("em", "i"):
                parse_inline(child, para, size, color, FONT_SERIF, bold=bold, italic=True)
            elif n == "a":
                parse_inline(child, para, size, YELLOW, font, bold=bold, italic=italic)
            elif n == "u":
                parse_inline(child, para, size, color, font, bold=bold, italic=italic)
            elif n == "br":
                r = para.add_run();  r.text = "\n"
            elif n not in ("style", "script", "li"):
                parse_inline(child, para, size, color, font, bold=bold, italic=italic)


def dash_list_box(slide, li_items, x, y, w, h, size=15):
    """
    Render a flat list of li-clone Tags as cyan em-dash items in one textbox.
    """
    tb = txbox(slide, x, y, w, h)
    tf = tb.text_frame
    for i, li in enumerate(li_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # Cyan dash
        r = p.add_run()
        r.text = "\u2014  "
        r.font.name = FONT_SANS;  r.font.size = Pt(size);  r.font.color.rgb = CYAN
        # Content
        parse_inline(li, p, size=size)
    return tb


def cyan_blockquote(slide, content_tag, x, y, w, bar_h=1.4, size=13):
    """Cyan left bar + body text — used for reading-discussion questions."""
    add_rect(slide, x, y, 0.18, bar_h, CYAN)
    tb = txbox(slide, x + 0.28, y, w - 0.3, bar_h)
    p = tb.text_frame.paragraphs[0]
    parse_inline(content_tag, p, size=size)
    return tb


# ═══════════════════════════════════════════════════════════════════════════════
# Slide classifer
# ═══════════════════════════════════════════════════════════════════════════════

def classify(section, index):
    classes = section.get("class", [])
    state   = section.get("data-state", "")
    if (section.find("h1") and index == 1) or state == "title":
        return "title"
    if "toc" in classes:
        return "toc"
    if "center" in classes:
        return "center"
    # Background-only slides (no content divs/lists) → center
    if section.get("data-background") and not section.find(["ul","ol","div"]):
        return "center"
    # .twothird with an <img> in its first column → reading
    tw = section.find(class_=re.compile(r"\btwothird\b"))
    if tw and tw.find("img"):
        return "reading"
    if section.find(class_=re.compile(r"\btwocol\b")):
        return "twocol"
    return "standard"


# ═══════════════════════════════════════════════════════════════════════════════
# Running-head extractor
# ═══════════════════════════════════════════════════════════════════════════════

def extract_running_heads(html_source):
    pat = re.compile(
        r'\.([\w-]+)\s+header:after\s*\{[^}]*content:\s*["\']([^"\']+)["\']',
        re.DOTALL)
    return {m.group(1): m.group(2) for m in pat.finditer(html_source)}


def get_running_head(section, heads):
    for token in section.get("class", []) + [section.get("data-state", "")]:
        if token in heads:
            return heads[token]
    return ""


# ═══════════════════════════════════════════════════════════════════════════════
# Slide renderers
# ═══════════════════════════════════════════════════════════════════════════════

def render_title(prs, sec, img_dir, head, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)
    bg = sec.get("data-background", "")
    if bg:
        set_image_bg(slide, resolve_img(bg, img_dir), transparency=25)

    h1 = sec.find("h1")
    if h1:
        tb = txbox(slide, 0.6, 0.5, 9.1, 2.9)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        parse_inline(h1, p, size=64, bold=True)

    h5 = sec.find("h5")
    if h5:
        simple_text(slide, 0.6, 3.55, 6.5, 0.9,
                    h5.get_text(" ", strip=True), size=22, bold=True)

    aside = sec.find("aside")
    if aside:
        simple_text(slide, 0.6, 4.55, 7, 0.4,
                    aside.get_text(" ", strip=True), size=14)

    add_running_head(slide, head)
    add_counter(slide, n, total)


def render_toc(prs, sec, head, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)

    h2 = sec.find("h2")
    simple_text(slide, 0.6, 0.1, 9, 1.15,
                h2.get_text(strip=True) if h2 else "Contents",
                font=FONT_SERIF, size=70, bold=True)

    # TOC items (nested-li chain)
    ol = sec.find("ol")
    items = []
    li = ol.find("li") if ol else None
    while li:
        u = li.find("u", recursive=False)
        text = u.get_text(strip=True) if u else "".join(
            str(c) for c in li.children if isinstance(c, NavigableString)).strip()
        if text:
            items.append(text)
        li = li.find("li", recursive=False)

    for i, text in enumerate(items):
        y = 1.35 + i * 0.85
        simple_text(slide, 0.4, y, 0.5, 0.62, f"{i+1}.",
                    size=26, align=PP_ALIGN.RIGHT)
        simple_text(slide, 1.0, y, 8.5, 0.62, text, size=28)
        add_rect(slide, 1.0, y + 0.62, min(len(text)*0.155+0.3, 8.8), 0.072, MAGENTA)

    aside = sec.find("aside")
    if aside:
        simple_text(slide, 0.6, 4.85, 5.5, 0.65,
                    aside.get_text(" ", strip=True), size=13)

    add_running_head(slide, head)
    add_counter(slide, n, total)


def render_center(prs, sec, img_dir, head, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)
    bg = sec.get("data-background", "")
    if bg:
        set_image_bg(slide, resolve_img(bg, img_dir), transparency=20)

    h1 = sec.find("h1");  h2 = sec.find("h2")
    heading = h1 or h2
    if heading:
        tb = txbox(slide, 0.5, 1.2, 9, 3.0)
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        parse_inline(heading, p, size=64 if h1 else 58,
                     font=FONT_SERIF, bold=True)

    y = 3.85
    for child in sec.children:
        if not isinstance(child, Tag):
            continue
        nm = (child.name or "").lower()
        if nm in ("h1","h2"):
            continue
        if nm == "p":
            simple_text(slide, 0.5, y, 9, 0.6,
                        child.get_text(" ", strip=True), size=18, align=PP_ALIGN.CENTER)
            y += 0.65
        elif nm == "aside":
            simple_text(slide, 0.5, y, 9, 0.6,
                        child.get_text(" ", strip=True), size=13,
                        color=MUTED, align=PP_ALIGN.CENTER)
            y += 0.65

    add_running_head(slide, head)
    add_counter(slide, n, total)


def render_standard(prs, sec, img_dir, head, n, total):
    """
    Standard slide: h2 (direct child of section) + content in twothird>div[0].
    Lists use nested <li> — unwrap with flatten_li().
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)

    # Heading (direct child of section, not inside twothird)
    h2 = sec.find("h2")
    if h2:
        simple_text(slide, 0.6, 0.3, 9, 0.85,
                    h2.get_text(strip=True), size=44, bold=True)

    # Content lives inside twothird > div[0]
    col = get_col0(sec)
    if col is None:
        # No twothird — fall back to section body
        col = sec

    cy = 1.28
    x, w = 0.6, 8.9

    for child in col.children:
        if not isinstance(child, Tag):
            continue
        nm = (child.name or "").lower()

        if nm == "h2":
            continue   # already rendered above

        elif nm in ("ul", "ol"):
            items = flatten_li(child)
            if items:
                box_h = min(len(items) * 0.6 + 0.15, 5.625 - cy - 0.2)
                dash_list_box(slide, items, x, cy, w, box_h, size=15)
                cy += box_h + 0.05

        elif nm == "p":
            text = child.get_text(" ", strip=True)
            if not text:
                continue
            # Render with inline formatting preserved
            tb = txbox(slide, x, cy, w, 0.65)
            p = tb.text_frame.paragraphs[0]
            parse_inline(child, p, size=15)
            cy += 0.7

        elif nm == "blockquote":
            cyan_blockquote(slide, child, x, cy, w, bar_h=1.2)
            cy += 1.35

        elif nm == "aside":
            simple_text(slide, x, cy, w, 0.55,
                        child.get_text(" ", strip=True), size=12, color=MUTED)
            cy += 0.6

        elif nm == "footer":
            simple_text(slide, x, 5.05, w, 0.35,
                        child.get_text(" ", strip=True), size=11, color=CYAN_L)

    add_running_head(slide, head)
    add_counter(slide, n, total)


def render_reading(prs, sec, img_dir, head, n, total):
    """
    Reading-discussion slide: portrait(s) on left, reading text on right.
    All content is in twothird > div[0].
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)

    col = get_col0(sec)
    imgs = col.find_all("img") if col else []

    # ── Portrait(s) ───────────────────────────────────────────────────────────
    if len(imgs) == 1:
        src = imgs[0].get("src","")
        add_image(slide, resolve_img(src, img_dir), 0.5, 0.35, 3.1, 3.1)
        # Italic label (p with style containing Happy Times / italic)
        for p_tag in col.find_all("p"):
            style = p_tag.get("style","")
            if "Happy Times" in style or "italic" in style:
                simple_text(slide, 0.5, 3.55, 3.1, 0.55,
                            p_tag.get_text(strip=True),
                            font=FONT_SERIF, size=20, italic=True)
                break
    elif len(imgs) >= 2:
        for j, img_tag in enumerate(imgs[:2]):
            add_image(slide, resolve_img(img_tag.get("src",""), img_dir),
                      0.5 + j*1.65, 0.35, 1.5, 1.9)
        for p_tag in col.find_all("p"):
            style = p_tag.get("style","")
            if "Happy Times" in style or "italic" in style:
                simple_text(slide, 0.5, 2.35, 3.3, 0.5,
                            p_tag.get_text(strip=True),
                            font=FONT_SERIF, size=20, italic=True)
                break

    # ── Right column: reading code + body text ─────────────────────────────
    rx, ry = 4.0, 0.35
    content_ps = [p for p in col.find_all("p")
                  if "Happy Times" not in p.get("style","")
                  and "italic" not in p.get("style","")]

    # The body text is one big <p> with <strong> for the code and <br> breaks.
    # We split it: first <strong> = reading code, rest = body + question
    for p_tag in content_ps:
        children = list(p_tag.children)

        # Extract reading code (first <strong>)
        first_strong = p_tag.find("strong")
        if first_strong:
            simple_text(slide, rx, ry, 5.7, 0.42,
                        first_strong.get_text(strip=True),
                        size=13, bold=True, color=CYAN)
            ry += 0.52

        # Find the split point: two consecutive <br> tags signal a new question
        # Build two segments: intro text, then question text
        seg1_nodes, seg2_nodes = [], []
        in_seg2 = False
        br_count = 0
        for child in children:
            if isinstance(child, Tag) and child.name == "strong":
                br_count = 0  # reset on any new tag
                if not in_seg2:
                    seg1_nodes.append(child)
                else:
                    seg2_nodes.append(child)
            elif isinstance(child, Tag) and child.name == "br":
                br_count += 1
                if br_count >= 2 and not in_seg2:
                    in_seg2 = True
                elif not in_seg2:
                    seg1_nodes.append(child)
                else:
                    seg2_nodes.append(child)
            else:
                br_count = 0
                if not in_seg2:
                    seg1_nodes.append(child)
                else:
                    seg2_nodes.append(child)

        # Render seg1 (body text)
        if seg1_nodes:
            # Make a fake tag to parse
            fake = BeautifulSoup("<span></span>", "html.parser").find("span")
            for node in seg1_nodes:
                fake.append(node.__copy__() if isinstance(node, Tag)
                            else NavigableString(str(node)))
            tb = txbox(slide, rx, ry, 5.7, 2.0)
            p = tb.text_frame.paragraphs[0]
            parse_inline(fake, p, size=13)
            ry += 2.1

        # Render seg2 as cyan-bar blockquote (the question)
        if seg2_nodes:
            fake2 = BeautifulSoup("<span></span>", "html.parser").find("span")
            for node in seg2_nodes:
                fake2.append(node.__copy__() if isinstance(node, Tag)
                             else NavigableString(str(node)))
            add_rect(slide, rx, ry, 0.18, 1.5, CYAN)
            tb2 = txbox(slide, rx+0.28, ry, 5.4, 1.5)
            p2 = tb2.text_frame.paragraphs[0]
            parse_inline(fake2, p2, size=13)

    add_running_head(slide, head)
    add_counter(slide, n, total)


def render_twocol(prs, sec, img_dir, head, n, total):
    """Two-column slide (practical exercise): h2 + dash list in each column."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_black_bg(slide)
    add_running_head(slide, head)

    specs = [{"x": 0.5, "w": 4.4}, {"x": 5.28, "w": 4.5}]
    for j, col in enumerate(get_cols(sec)[:2]):
        x, w = specs[j]["x"], specs[j]["w"]
        cy = 0.38

        h2 = col.find("h2")
        if h2:
            simple_text(slide, x, cy, w, 0.65,
                        h2.get_text(strip=True), size=30, bold=True)
            cy += 0.72

        for list_tag in col.find_all(["ul","ol"], recursive=False):
            items = flatten_li(list_tag)
            if items:
                box_h = 5.625 - cy - 0.25
                dash_list_box(slide, items, x, cy, w, box_h, size=13)
                cy += box_h

    # Divider
    add_rect(slide, 5.06, 0.3, 0.02, 5.1, MUTED)
    add_counter(slide, n, total)


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def convert(html_path: Path, out_path: Path):
    html_source = html_path.read_text(encoding="utf-8")
    soup        = BeautifulSoup(html_source, "html.parser")
    img_dir     = html_path.parent / "img"
    heads       = extract_running_heads(html_source)

    slides_div = soup.find("div", class_="slides")
    if not slides_div:
        sys.exit("ERROR: <div class='slides'> not found")
    wrapper    = slides_div.find("section")
    slide_tags = (wrapper.find_all("section", recursive=False) if wrapper
                  else slides_div.find_all("section", recursive=False))
    total = len(slide_tags)

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for i, sec in enumerate(slide_tags, 1):
        kind = classify(sec, i)
        head = get_running_head(sec, heads)
        print(f"  [{i:02d}/{total}]  {kind:<10}  head={head!r}")

        if   kind == "title":   render_title(prs, sec, img_dir, head, i, total)
        elif kind == "toc":     render_toc(prs, sec, head, i, total)
        elif kind == "center":  render_center(prs, sec, img_dir, head, i, total)
        elif kind == "reading": render_reading(prs, sec, img_dir, head, i, total)
        elif kind == "twocol":  render_twocol(prs, sec, img_dir, head, i, total)
        else:                   render_standard(prs, sec, img_dir, head, i, total)

    prs.save(str(out_path))
    print(f"\n✓  Saved → {out_path}")


def main():
    ap = argparse.ArgumentParser(
        description="Convert reveal.js HTML (publictimes_CMYK design) → .pptx")
    ap.add_argument("html", help="Path to the reveal.js HTML file")
    ap.add_argument("--out", "-o", help="Output .pptx (default: same name as HTML)")
    args = ap.parse_args()

    html_path = Path(args.html).resolve()
    if not html_path.exists():
        sys.exit(f"ERROR: {html_path} not found")
    out_path = Path(args.out).resolve() if args.out else html_path.with_suffix(".pptx")

    print(f"Input:  {html_path}")
    print(f"Output: {out_path}\n")
    convert(html_path, out_path)


if __name__ == "__main__":
    main()
